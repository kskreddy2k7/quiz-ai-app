import io
from pathlib import Path

from fastapi import UploadFile

from app.core.config import get_settings
from app.core.logging import logger

SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".jpg", ".jpeg", ".png"}


def validate_extension(filename: str) -> None:
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            "Unsupported file type. Please upload PDF, PPT/PPTX, DOC/DOCX, or image files."
        )


def read_upload(upload: UploadFile) -> str:
    validate_extension(upload.filename or "")
    suffix = Path(upload.filename or "").suffix.lower()
    data = upload.file.read()
    settings = get_settings()
    max_bytes = settings.max_upload_mb * 1024 * 1024
    if len(data) > max_bytes:
        raise ValueError(f"File too large. Max size is {settings.max_upload_mb}MB.")

    if suffix == ".txt":
        return data.decode("utf-8", errors="ignore")
    if suffix == ".pdf":
        return _read_pdf(data)
    if suffix in {".jpg", ".jpeg", ".png"}:
        return _read_image(data)
    if suffix in {".doc", ".ppt"}:
        return _read_legacy(data, suffix)
    if suffix == ".docx":
        return _read_docx(data)
    if suffix == ".pptx":
        return _read_pptx(data)

    raise ValueError("Unsupported file format.")


def _read_pdf(data: bytes) -> str:
    try:
        import pdfplumber

        with pdfplumber.open(io.BytesIO(data)) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)
    except Exception as exc:  # pragma: no cover - runtime dependency
        logger.warning("PDF extraction failed: %s", exc)
        raise ValueError("PDF text extraction failed. Please try a different file.")


def _read_docx(data: bytes) -> str:
    try:
        from docx import Document

        doc = Document(io.BytesIO(data))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    except Exception as exc:  # pragma: no cover
        logger.warning("DOCX extraction failed: %s", exc)
        raise ValueError("DOCX text extraction failed. Please try a different file.")


def _read_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation

        ppt = Presentation(io.BytesIO(data))
        slides = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides.append(shape.text)
        return "\n".join(slides)
    except Exception as exc:  # pragma: no cover
        logger.warning("PPTX extraction failed: %s", exc)
        raise ValueError("PPTX text extraction failed. Please try a different file.")


def _read_legacy(data: bytes, suffix: str) -> str:
    try:
        import tempfile

        import textract

        with tempfile.NamedTemporaryFile(suffix=suffix) as temp_file:
            temp_file.write(data)
            temp_file.flush()
            text = textract.process(temp_file.name)
        return text.decode("utf-8", errors="ignore")
    except Exception as exc:  # pragma: no cover
        logger.warning("Legacy %s extraction failed: %s", suffix, exc)
        raise ValueError(
            f"{suffix.upper()} extraction failed. Please upload a {suffix}x file if possible."
        )


def _read_image(data: bytes) -> str:
    try:
        from PIL import Image
        import pytesseract

        image = Image.open(io.BytesIO(data))
        return pytesseract.image_to_string(image)
    except Exception as exc:  # pragma: no cover
        logger.warning("Image OCR failed: %s", exc)
        raise ValueError("Image text extraction failed. Please try a clearer image.")
