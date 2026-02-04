from pathlib import Path
import fitz
from docx import Document
from pptx import Presentation


def extract_text(path: Path) -> str:
    if path.suffix == ".pdf":
        doc = fitz.open(path)
        return "\n".join(p.get_text() for p in doc)

    if path.suffix in (".doc", ".docx"):
        return "\n".join(p.text for p in Document(path).paragraphs)

    if path.suffix in (".ppt", ".pptx"):
        prs = Presentation(path)
        return "\n".join(
            s.text for slide in prs.slides for s in slide.shapes if hasattr(s, "text")
        )
    return ""
