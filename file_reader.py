def read_file(path):
    text = ""

    if path.endswith(".pdf"):
        import fitz
        pdf = fitz.open(path)
        for page in pdf:
            text += page.get_text()

    elif path.endswith(".docx"):
        from docx import Document
        doc = Document(path)
        text = " ".join(p.text for p in doc.paragraphs)

    elif path.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "

    elif path.endswith(".txt"):
        with open(path, "r", encoding="utf-8") as f:
            text = f.read()

    return textdef read_file(path):
    text = ""

    if path.endswith(".pdf"):
        import fitz
        pdf = fitz.open(path)
        for page in pdf:
            text += page.get_text()

    elif path.endswith(".docx"):
        from docx import Document
        doc = Document(path)
        text = " ".join(p.text for p in doc.paragraphs)

    elif path.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "

    elif path.endswith(".txt"):
        with open(path, "r", encoding="utf-8") as f:
            text = f.read()

    return text
